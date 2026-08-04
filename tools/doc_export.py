"""
文档导出 —— Markdown → DOCX / MD 字节流生成（零新增依赖）。

- 管线：markdown 库（tables/fenced_code 扩展）→ HTML → BeautifulSoup 遍历
  → python-docx 渲染，全部为项目既有依赖。
- 覆盖语法：h1-h6 / 段落 / 加粗 / 斜体 / 删除线 / 行内代码 / 链接 / 图片占位
  / 有序无序嵌套列表 / 围栏代码块（灰底等宽）/ 引用块 / 表格 / 分隔线。
- CPU 密集（大文档解析+构建 XML），调用方必须经 asyncio.to_thread 执行，
  禁止在事件循环上直接调用（对话零阻塞架构铁律）。
"""
from __future__ import annotations

import io
import logging
import re

logger = logging.getLogger("second_person.doc_export")

_CODE_FONT = "Consolas"
_CODE_BG = "F2F3F5"      # 代码块底色
_QUOTE_COLOR = "666666"  # 引用文字灰
_LINK_COLOR = "2563EB"   # 链接蓝


def sanitize_filename(name: str, default: str = "导出文档") -> str:
    """去除文件名非法字符，兜底默认名。"""
    name = re.sub(r'[\\/:*?"<>|\r\n]+', "_", (name or "").strip())
    return (name[:80] or default)


def md_to_docx_bytes(md_text: str, title: str | None = None) -> bytes:
    """Markdown 文本 → docx 文件字节。CPU 密集，须在工作线程调用。"""
    import markdown
    from bs4 import BeautifulSoup
    from docx import Document

    html = markdown.markdown(
        md_text or "", extensions=["tables", "fenced_code", "sane_lists"])
    soup = BeautifulSoup(html, "html.parser")

    doc = Document()
    _setup_base_style(doc)
    if title:
        doc.add_heading(title, level=0)
    for el in soup.children:
        _render_block(doc, el)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ---- 文档级样式 -----------------------------------------------------------
def _setup_base_style(doc) -> None:
    from docx.oxml.ns import qn
    from docx.shared import Pt
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)
    # 中文字体：eastAsia 需走 rPr 底层属性
    style.element.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")


# ---- 块级渲染 -------------------------------------------------------------
def _render_block(doc, el, quote_depth: int = 0) -> None:
    name = getattr(el, "name", None)
    if name is None:  # 顶层游离文本（通常是空白）
        text = str(el).strip()
        if text:
            _fill_inline(doc.add_paragraph(), [el])
        return
    if name in ("h1", "h2", "h3", "h4", "h5", "h6"):
        doc.add_heading(el.get_text(strip=True), level=int(name[1]))
    elif name == "p":
        p = doc.add_paragraph()
        if quote_depth:
            _style_quote(p, quote_depth)
        _fill_inline(p, el.children, quote=bool(quote_depth))
    elif name in ("ul", "ol"):
        _render_list(doc, el, ordered=(name == "ol"), level=0)
    elif name == "pre":
        code = el.find("code")
        _render_code_block(doc, (code or el).get_text())
    elif name == "blockquote":
        for child in el.children:
            if getattr(child, "name", None):
                _render_block(doc, child, quote_depth=quote_depth + 1)
    elif name == "table":
        _render_table(doc, el)
    elif name == "hr":
        _render_hr(doc)
    else:  # 未覆盖标签（div 等）：降级为纯文本段落
        text = el.get_text(strip=True)
        if text:
            doc.add_paragraph(text)


def _render_list(doc, el, ordered: bool, level: int) -> None:
    # python-docx 内置样式仅到 3 级，更深层夹紧到 3
    suffix = "" if level == 0 else f" {min(level + 1, 3)}"
    style = ("List Number" if ordered else "List Bullet") + suffix
    for li in el.find_all("li", recursive=False):
        p = doc.add_paragraph(style=style)
        nested = []
        inline = []
        for child in li.children:
            if getattr(child, "name", None) in ("ul", "ol"):
                nested.append(child)
            elif getattr(child, "name", None) == "p":
                # 宽松列表：li 内含 p，取其行内内容
                inline.extend(list(child.children))
            else:
                inline.append(child)
        _fill_inline(p, inline)
        for sub in nested:
            _render_list(doc, sub, ordered=(sub.name == "ol"), level=level + 1)


def _render_code_block(doc, code_text: str) -> None:
    from docx.shared import Pt
    p = doc.add_paragraph()
    _shade_paragraph(p, _CODE_BG)
    lines = code_text.rstrip("\n").split("\n")
    for i, line in enumerate(lines):
        run = p.add_run(line)
        run.font.name = _CODE_FONT
        run.font.size = Pt(9)
        if i < len(lines) - 1:
            run.add_break()


def _render_table(doc, el) -> None:
    rows = el.find_all("tr")
    if not rows:
        return
    ncols = max(len(r.find_all(["th", "td"])) for r in rows)
    if not ncols:
        return
    table = doc.add_table(rows=len(rows), cols=ncols)
    table.style = "Table Grid"
    for ri, tr in enumerate(rows):
        for ci, cell in enumerate(tr.find_all(["th", "td"])[:ncols]):
            para = table.cell(ri, ci).paragraphs[0]
            _fill_inline(para, cell.children, bold_all=(cell.name == "th"))


def _render_hr(doc) -> None:
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    p = doc.add_paragraph()
    pPr = p._p.get_or_add_pPr()
    borders = OxmlElement("w:pBdr")
    bottom = OxmlElement("w:bottom")
    bottom.set(qn("w:val"), "single")
    bottom.set(qn("w:sz"), "6")
    bottom.set(qn("w:color"), "CCCCCC")
    borders.append(bottom)
    pPr.append(borders)


# ---- 行内渲染 -------------------------------------------------------------
def _fill_inline(p, nodes, bold=False, italic=False, strike=False,
                 code=False, quote=False, bold_all=False) -> None:
    from docx.shared import Pt, RGBColor
    for node in nodes:
        name = getattr(node, "name", None)
        if name is None:
            text = str(node)
            if not text:
                continue
            run = p.add_run(text)
            run.bold = bold or bold_all
            run.italic = italic
            run.font.strike = strike
            if code:
                run.font.name = _CODE_FONT
                run.font.size = Pt(9)
                _shade_run(run, _CODE_BG)
            if quote:
                run.font.color.rgb = RGBColor.from_string(_QUOTE_COLOR)
        elif name in ("strong", "b"):
            _fill_inline(p, node.children, True, italic, strike, code,
                         quote, bold_all)
        elif name in ("em", "i"):
            _fill_inline(p, node.children, bold, True, strike, code,
                         quote, bold_all)
        elif name in ("del", "s", "strike"):
            _fill_inline(p, node.children, bold, italic, True, code,
                         quote, bold_all)
        elif name == "code":
            _fill_inline(p, node.children, bold, italic, strike, True,
                         quote, bold_all)
        elif name == "a":
            _add_hyperlink(p, node.get("href", ""),
                           node.get_text() or node.get("href", ""))
        elif name == "img":
            alt = node.get("alt") or node.get("src", "")
            p.add_run(f"[图片: {alt}]").italic = True
        elif name == "br":
            if p.runs:
                p.runs[-1].add_break()
        else:  # span 等其余行内容器
            _fill_inline(p, node.children, bold, italic, strike, code,
                         quote, bold_all)


def _add_hyperlink(p, url: str, text: str) -> None:
    """python-docx 无内置超链接 API，走 relationship + w:hyperlink 底层构建。"""
    from docx.opc.constants import RELATIONSHIP_TYPE as RT
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    try:
        r_id = p.part.relate_to(url, RT.HYPERLINK, is_external=True)
    except Exception:  # noqa: BLE001 - 非法 URL 降级为普通文本
        p.add_run(text)
        return
    link = OxmlElement("w:hyperlink")
    link.set(qn("r:id"), r_id)
    r = OxmlElement("w:r")
    rPr = OxmlElement("w:rPr")
    color = OxmlElement("w:color")
    color.set(qn("w:val"), _LINK_COLOR)
    u = OxmlElement("w:u")
    u.set(qn("w:val"), "single")
    rPr.append(color)
    rPr.append(u)
    r.append(rPr)
    t = OxmlElement("w:t")
    t.text = text
    t.set(qn("xml:space"), "preserve")
    r.append(t)
    link.append(r)
    p._p.append(link)


# ---- 底层着色 -------------------------------------------------------------
def _shade_paragraph(p, hex_color: str) -> None:
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:fill"), hex_color)
    p._p.get_or_add_pPr().append(shd)


def _shade_run(run, hex_color: str) -> None:
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:fill"), hex_color)
    run._r.get_or_add_rPr().append(shd)


def _style_quote(p, depth: int) -> None:
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.shared import Pt
    p.paragraph_format.left_indent = Pt(14 * depth)
    # 左侧竖线：引用视觉标记
    pPr = p._p.get_or_add_pPr()
    borders = OxmlElement("w:pBdr")
    left = OxmlElement("w:left")
    left.set(qn("w:val"), "single")
    left.set(qn("w:sz"), "12")
    left.set(qn("w:color"), "D1D5DB")
    borders.append(left)
    pPr.append(borders)


# ---- PPTX 导出（Markdown → 演示文稿）---------------------------------------
# 页面组织：首个 h1 为封面标题；每个 h2 开启新页；h3+ 为页内小节；
# 无 h2 时全文单页；页数 / 页字符超限截断并标注（防超大文档拖垮渲染）。
# CPU 密集，调用方必须经 asyncio.to_thread 执行（对话零阻塞铁律）。
_PPT_SLIDE_W = 13.333          # 16:9 宽（英寸）
_PPT_SLIDE_H = 7.5
_PPT_MAX_PAGES = 40
_PPT_MAX_PAGE_CHARS = 800
_PPT_MAX_TABLE_COLS = 6
_PPT_MAX_TABLE_ROWS = 18
_PPT_MAX_CODE_LINES = 20
_PPT_CELL_TEXT_MAX = 40
_PPT_CODE_BG = "F2F3F5"
_PPT_CODE_COLOR = "333333"
_PPT_QUOTE_COLOR = "666666"
_PPT_BODY_BOTTOM = 7.2         # 内容页可写底部边界（英寸）
_PPT_FONT_CN = "微软雅黑"


def _set_pptx_font(run, name, size=None, bold=None, italic=None, color=None):
    """设置 run 字体：latin + eastAsia 双通道（中文渲染必需）。"""
    from pptx.oxml.ns import qn
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", name)
    if size is not None:
        run.font.size = size
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color


def md_to_pptx_bytes(md_text: str, title: str | None = None) -> bytes:
    """Markdown 文本 → pptx 文件字节。CPU 密集，须在工作线程调用。

    python-pptx 未安装时抛 RuntimeError（中文提示），由 generate_document 兜底。
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        raise RuntimeError(
            "PPT 导出功能不可用：缺少 python-pptx 依赖，安装后重试") from None

    import markdown
    from bs4 import BeautifulSoup

    html = markdown.markdown(
        md_text or "", extensions=["tables", "fenced_code", "sane_lists"])
    soup = BeautifulSoup(html, "html.parser")

    prs = Presentation()
    prs.slide_width = Inches(_PPT_SLIDE_W)
    prs.slide_height = Inches(_PPT_SLIDE_H)

    blocks = [b for b in soup.children if getattr(b, "name", None)]
    h1_text = next((b.get_text(strip=True) for b in blocks if b.name == "h1"), "")
    _add_pptx_cover(prs, title or h1_text or "文档",
                    h1_text if title and h1_text else None)

    # 内容页状态：y=页内游标，chars=页内已写字符数（防单页超载）
    st = {"slide": None, "y": 0.35, "chars": 0, "pages": 1, "truncated": False}
    h1_seen = False  # 首个 h1 已用作封面；后续 h1 降级为 h2（新页）防丢内容

    def _new_page(page_title: str = "") -> bool:
        if st["pages"] >= _PPT_MAX_PAGES:
            st["truncated"] = True
            return False
        st["pages"] += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        st["slide"], st["y"], st["chars"] = slide, 0.35, 0
        if page_title:
            _add_pptx_heading(slide, page_title, size=24, y=0.2)
            st["y"] += 0.55
            st["chars"] += len(page_title)
        return True

    def _place(height: float, chars: int) -> bool:
        """检查空间：放不下则换页；返回 False 表示已超页数上限截断。"""
        if st["slide"] is not None and st["y"] + height <= _PPT_BODY_BOTTOM \
                and st["chars"] + chars <= _PPT_MAX_PAGE_CHARS:
            return True
        return _new_page()

    def _advance(height: float, chars: int) -> None:
        st["y"] += height
        st["chars"] += chars

    for b in blocks:
        if st["truncated"]:
            break
        name = b.name
        if name == "h1":
            if not h1_seen:
                h1_seen = True
                continue  # 首个 h1 已用作封面
            if not _new_page(b.get_text(strip=True)):
                break
            continue
        if name == "h2":
            if not _new_page(b.get_text(strip=True)):
                break
            continue
        if st["slide"] is None and not _new_page():
            break
        if name in ("h3", "h4", "h5", "h6"):
            text = b.get_text(strip=True)
            if not text:
                continue
            if not _place(0.45, len(text) + 10):
                break
            _add_pptx_heading(st["slide"], text, size=18)
            _advance(0.45, len(text) + 10)
        elif name == "table":
            rows = b.find_all("tr")
            ncols = max((len(r.find_all(["th", "td"])) for r in rows), default=0)
            if not rows or not ncols:
                continue
            ncols = min(ncols, _PPT_MAX_TABLE_COLS)
            # 行数截断防整表丢弃：超限只渲染前 N 行并标注
            clipped_rows = rows[:_PPT_MAX_TABLE_ROWS]
            h = 0.35 + len(clipped_rows) * 0.3
            if not _place(h, 40):
                break
            _add_pptx_table(st["slide"], clipped_rows, ncols, st["y"], h)
            if len(rows) > _PPT_MAX_TABLE_ROWS:
                _add_pptx_note(st["slide"],
                               f"（表格已截断，共 {len(rows)} 行仅显示前 {_PPT_MAX_TABLE_ROWS} 行）")
                _advance(0.3, 40)
            _advance(h, 40)
        elif name == "pre":
            code = (b.find("code") or b).get_text()
            lines = code.rstrip("\n").split("\n")
            clipped = lines[:_PPT_MAX_CODE_LINES]
            h = 0.28 + len(clipped) * 0.19
            if not _place(h, min(len(code), 400)):
                break
            _add_pptx_code(st["slide"], clipped, st["y"],
                           "（代码已截断）" if len(lines) > _PPT_MAX_CODE_LINES else "")
            _advance(h, min(len(code), 400))
        elif name == "blockquote":
            text = b.get_text(" ", strip=True)
            if not text:
                continue
            if not _place(0.32, len(text)):
                break
            _add_pptx_quote(st["slide"], text, st["y"])
            _advance(0.32, len(text))
        elif name in ("ul", "ol"):
            items: list[tuple[str, int, bool]] = []
            _collect_pptx_items(b, items)
            if not items:
                continue
            h = 0.32 * len(items)
            if not _place(h, sum(len(i[0]) for i in items)):
                break
            _add_pptx_list(st["slide"], items, st["y"])
            _advance(h, sum(len(i[0]) for i in items))
        else:  # p 及其余：段落
            text = b.get_text(" ", strip=True)
            if not text:
                continue
            if not _place(0.32, len(text)):
                break
            _add_pptx_para(st["slide"], b, st["y"])
            _advance(0.32, len(text))

    if st["truncated"] and st["slide"] is not None:
        _add_pptx_note(st["slide"], "（内容已截断，超出页数上限）")

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _add_pptx_cover(prs, title: str, subtitle: str | None) -> None:
    from pptx.util import Pt
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
    if slide.shapes.title:
        slide.shapes.title.text = title
        for p in slide.shapes.title.text_frame.paragraphs:
            for run in p.runs:
                _set_pptx_font(run, _PPT_FONT_CN, size=Pt(36), bold=True)
    if subtitle and len(slide.placeholders) > 1:
        ph = slide.placeholders[1]
        ph.text = subtitle
        for p in ph.text_frame.paragraphs:
            for run in p.runs:
                _set_pptx_font(run, _PPT_FONT_CN, size=Pt(16))


def _add_pptx_textbox(slide, y: float, height: float = 0.35):
    from pptx.util import Inches
    tb = slide.shapes.add_textbox(
        Inches(0.55), Inches(y), Inches(_PPT_SLIDE_W - 1.1), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def _add_pptx_heading(slide, text: str, size: int = 24, y: float = 0.2) -> None:
    from pptx.util import Inches, Pt
    tb = slide.shapes.add_textbox(
        Inches(0.55), Inches(y), Inches(_PPT_SLIDE_W - 1.1), Inches(0.55))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_pptx_font(run, _PPT_FONT_CN, size=Pt(size), bold=True)


def _fill_pptx_runs(p, nodes, bold=False, italic=False, strike=False,
                    code=False) -> None:
    """行内渲染：加粗/斜体/删除线/行内代码；链接降级为普通文本。"""
    from pptx.dml.color import RGBColor
    from pptx.util import Pt
    for node in nodes:
        name = getattr(node, "name", None)
        if name is None:
            text = str(node)
            if not text:
                continue
            run = p.add_run()
            run.text = text.replace("\n", "\v")
            _set_pptx_font(
                run, "Consolas" if code else _PPT_FONT_CN,
                size=Pt(9) if code else None, bold=bold, italic=italic,
                color=RGBColor.from_string(_PPT_CODE_COLOR) if code else None)
            if strike:
                run._r.get_or_add_rPr().set("strike", "sngStrike")
        elif name in ("strong", "b"):
            _fill_pptx_runs(p, node.children, True, italic, strike, code)
        elif name in ("em", "i"):
            _fill_pptx_runs(p, node.children, bold, True, strike, code)
        elif name in ("del", "s", "strike"):
            _fill_pptx_runs(p, node.children, bold, italic, True, code)
        elif name == "code":
            _fill_pptx_runs(p, node.children, bold, italic, strike, True)
        elif name == "a":
            run = p.add_run()
            run.text = (node.get_text() or node.get("href", "")).replace("\n", "\v")
            _set_pptx_font(run, _PPT_FONT_CN)
        elif name == "img":
            run = p.add_run()
            run.text = f"[图片: {node.get('alt') or node.get('src', '')}]"
            _set_pptx_font(run, _PPT_FONT_CN)
        elif name == "br":
            run = p.add_run()
            run.text = "\v"
            _set_pptx_font(run, _PPT_FONT_CN)
        else:  # span 等其余行内容器
            _fill_pptx_runs(p, node.children, bold, italic, strike, code)


def _add_pptx_para(slide, el, y: float) -> None:
    tb, tf = _add_pptx_textbox(slide, y)
    _fill_pptx_runs(tf.paragraphs[0], el.children)


def _collect_pptx_items(el, out: list, level: int = 0,
                        ordered: bool = False) -> None:
    """收集列表项 (文本, 缩进级别, 是否有序)；嵌套 ≤2 级夹紧。"""
    for li in el.find_all("li", recursive=False):
        parts: list[str] = []
        nested = []
        for child in li.children:
            if getattr(child, "name", None) in ("ul", "ol"):
                nested.append(child)
            elif getattr(child, "name", None) == "p":
                parts.append(child.get_text(" ", strip=True))
            elif getattr(child, "name", None) is not None:
                parts.append(child.get_text(" ", strip=True))
            else:
                parts.append(str(child).strip())
        text = " ".join(x for x in parts if x).strip()
        if text:
            out.append((text, min(level, 2), ordered))
        for sub in nested:
            _collect_pptx_items(sub, out, level + 1, ordered=(sub.name == "ol"))


def _add_pptx_list(slide, items: list, y: float) -> None:
    from pptx.util import Pt
    tb, tf = _add_pptx_textbox(slide, y, height=0.32 * max(len(items), 1))
    for i, (text, level, ordered) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        run = p.add_run()
        run.text = ("1. " if ordered else "• ") + text
        _set_pptx_font(run, _PPT_FONT_CN, size=Pt(14))


def _add_pptx_code(slide, lines: list[str], y: float, note: str = "") -> None:
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    h = 0.28 + len(lines) * 0.19 + (0.2 if note else 0)
    tb = slide.shapes.add_textbox(
        Inches(0.55), Inches(y), Inches(_PPT_SLIDE_W - 1.1), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.08)
    tb.fill.solid()
    tb.fill.fore_color.rgb = RGBColor.from_string(_PPT_CODE_BG)
    run = tf.paragraphs[0].add_run()
    run.text = "\v".join(lines)
    _set_pptx_font(run, "Consolas", size=Pt(11),
                   color=RGBColor.from_string(_PPT_CODE_COLOR))
    if note:
        r2 = tf.add_paragraph().add_run()
        r2.text = note
        _set_pptx_font(r2, _PPT_FONT_CN, size=Pt(9),
                       color=RGBColor.from_string(_PPT_QUOTE_COLOR))


def _add_pptx_quote(slide, text: str, y: float) -> None:
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    tb, tf = _add_pptx_textbox(slide, y)
    run = tf.paragraphs[0].add_run()
    run.text = text
    _set_pptx_font(run, _PPT_FONT_CN, size=Pt(13), italic=True,
                   color=RGBColor.from_string(_PPT_QUOTE_COLOR))


def _add_pptx_table(slide, rows, ncols: int, y: float, height: float) -> None:
    from pptx.util import Inches, Pt
    gf = slide.shapes.add_table(
        len(rows), ncols, Inches(0.55), Inches(y),
        Inches(_PPT_SLIDE_W - 1.1), Inches(max(height, 0.4)))
    table = gf.table
    col_w = (_PPT_SLIDE_W - 1.1) / ncols
    for ci in range(ncols):
        table.columns[ci].width = Inches(col_w)
    for ri, tr in enumerate(rows):
        cells = tr.find_all(["th", "td"])[:ncols]
        for ci in range(ncols):
            cell = table.cell(ri, ci)
            text = cells[ci].get_text(" ", strip=True) if ci < len(cells) else ""
            cell.text = text[:_PPT_CELL_TEXT_MAX]
            is_head = ci < len(cells) and cells[ci].name == "th"
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    _set_pptx_font(run, _PPT_FONT_CN, size=Pt(11),
                                   bold=is_head)


def _add_pptx_note(slide, text: str) -> None:
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    tb = slide.shapes.add_textbox(
        Inches(0.55), Inches(_PPT_BODY_BOTTOM - 0.4),
        Inches(_PPT_SLIDE_W - 1.1), Inches(0.35))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = text
    _set_pptx_font(run, _PPT_FONT_CN, size=Pt(10),
                   color=RGBColor.from_string(_PPT_QUOTE_COLOR))


# ---- XLSX 导出（Markdown → Excel 工作簿）-----------------------------------
# 每个 markdown 表格一个 sheet；表格间文本归入「内容」sheet；
# 公式注入防护：= + - @ 开头单元格前置单引号。
# CPU 密集，调用方必须经 asyncio.to_thread 执行（对话零阻塞铁律）。
_XLSX_MAX_ROWS = 10000
_XLSX_MAX_CELL_CHARS = 32767
_XLSX_MAX_SHEET_NAME = 31
_XLSX_HEADER_BG = "F2F3F5"


def md_to_xlsx_bytes(md_text: str, title: str | None = None) -> bytes:
    """Markdown 文本 → xlsx 文件字节。CPU 密集，须在工作线程调用。

    openpyxl 未安装时抛 RuntimeError（中文提示），由 generate_document 兜底。
    """
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill
        from openpyxl.utils import get_column_letter
    except ImportError:
        raise RuntimeError(
            "Excel 导出功能不可用：缺少 openpyxl 依赖，安装后重试") from None

    import markdown
    from bs4 import BeautifulSoup

    html = markdown.markdown(
        md_text or "", extensions=["tables", "fenced_code", "sane_lists"])
    soup = BeautifulSoup(html, "html.parser")

    wb = Workbook()
    wb.remove(wb.active)
    st: dict = {"seq": 0, "pending_title": None, "content": None,
                "title_written": False}

    def _cell(raw) -> str:
        text = str(raw or "").strip()
        if len(text) > _XLSX_MAX_CELL_CHARS:
            text = text[:_XLSX_MAX_CELL_CHARS]
        if text and text[0] in "=+-@":
            return "'" + text  # 公式注入防护
        return text

    def _sheet_name(base: str | None) -> str:
        import re as _re
        name = _re.sub(r"[\[\]:*?/\\]", "", (base or "").strip())
        name = (name or "表格")[:_XLSX_MAX_SHEET_NAME]
        st["seq"] += 1
        candidate = name if st["seq"] == 1 else f"{name}({st['seq']})"
        while candidate in wb.sheetnames:
            st["seq"] += 1
            candidate = f"{name}({st['seq']})"
        return candidate

    def _write_row(sheet, values: list, header: bool = False) -> bool:
        """写一行；超行数上限时标注一次并返回 False（停止写入该 sheet）。"""
        if sheet.max_row >= _XLSX_MAX_ROWS:
            key = f"trunc_{sheet.title}"
            if not st.get(key):
                st[key] = True
                sheet.append(["（内容已截断，超出行数上限）"])
            return False
        sheet.append([_cell(v) for v in values])
        if header:
            for cell in sheet[sheet.max_row]:
                cell.font = Font(bold=True)
                cell.fill = PatternFill("solid", fgColor=_XLSX_HEADER_BG)
        return True

    def _content_sheet():
        if st["content"] is None:
            st["content"] = wb.create_sheet("内容")
            if title and not st["title_written"]:
                _write_row(st["content"], [title], header=True)
                st["title_written"] = True
        return st["content"]

    for b in soup.children:
        name = getattr(b, "name", None)
        if name == "h1":
            text = b.get_text(strip=True)
            cs = _content_sheet()  # 先触发 title 写入，避免与 h1 重复
            if text:
                if st["title_written"] and cs.max_row > 1 \
                        and cs["A1"].value != text:
                    # title 已写但与 h1 不同：h1 降级为普通行，不丢内容
                    _write_row(cs, [text])
                elif not st["title_written"]:
                    _write_row(cs, [text], header=True)
                    st["title_written"] = True
                st["pending_title"] = text
        elif name in ("h2", "h3", "h4", "h5", "h6"):
            text = b.get_text(strip=True)
            if text:
                st["pending_title"] = text
        elif name == "table":
            rows = b.find_all("tr")
            if not rows:
                continue
            ncols = max((len(r.find_all(["th", "td"])) for r in rows), default=0)
            if not ncols:
                continue
            sheet = wb.create_sheet(_sheet_name(st["pending_title"]))
            for ri, tr in enumerate(rows):
                cells = tr.find_all(["th", "td"])
                values = [cells[ci].get_text(" ", strip=True)
                          if ci < len(cells) else "" for ci in range(ncols)]
                if not _write_row(sheet, values, header=(ri == 0)):
                    break
        elif name == "pre":
            code = (b.find("code") or b).get_text()
            for line in code.rstrip("\n").split("\n"):
                _write_row(_content_sheet(), ["[代码块]", line])
        elif name in ("p", "ul", "ol", "blockquote"):
            text = b.get_text(" ", strip=True)
            if text:
                _write_row(_content_sheet(), [text])
        # 其余标签忽略

    if not wb.sheetnames:
        ws = wb.create_sheet("内容")
        _write_row(ws, [title or "（空文档）"], header=bool(title))

    # 列宽自适应（表头 + 数据最长值）
    for ws in wb.worksheets:
        for col in ws.columns:
            max_len = max((len(str(c.value or "")) for c in col), default=0)
            letter = get_column_letter(col[0].column)
            ws.column_dimensions[letter].width = min(max(max_len * 1.2, 8), 40)

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()
