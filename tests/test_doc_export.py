"""文档导出多格式测试（generate_document: docx / md / pptx / xlsx）。

覆盖：
1. 转换器全语法（标题/列表/表格/代码块/引用/行内样式）与文件可重新打开
2. PPTX 边界：页数上限 / 代码块截断 / 大表格截断 / 多 h1 降级 / 空文档
3. XLSX 边界：公式注入防护 / sheet 名清洗冲突 / 标题去重 / 空文档
4. 工具级：四格式生成与落盘、别名归一化、非法格式与空内容拒绝、schema 枚举
5. 前端无 docx/md 格式硬编码（新格式对 UI 零改动回归防线）
运行：python tests/test_doc_export.py（退出码 0 = 全部通过）
"""
import asyncio
import io
import sys
import tempfile
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

FAIL: list[str] = []


def check(name: str, cond: bool, detail: str = "") -> None:
    print(f"[{'PASS' if cond else 'FAIL'}] {name}" + (f" - {detail}" if detail else ""))
    if not cond:
        FAIL.append(name)


FULL_MD = """# 季度报告

## 一、业绩概览

**要点**：本季度营收 *增长* 12%，~~旧数据~~已作废。

- 第一项
- 第二项
  - 嵌套子项
- 第三项

1. 有序甲
2. 有序乙

> 引用：管理层认为市场回暖。

| 指标 | Q1 | Q2 |
| --- | --- | --- |
| 营收 | 100 | 112 |
| 利润 | 20 | 25 |

## 二、技术细节

代码示例：

```python
def hello(name):
    print(f"hi {name}")
```

行内 `code` 与[链接](https://example.com)与![图](x.png)。

### 小节标题

结尾段落。
"""


def test_converters() -> None:
    from tools.doc_export import (md_to_docx_bytes, md_to_pptx_bytes,
                                  md_to_xlsx_bytes)
    for name, fn in (("docx", md_to_docx_bytes), ("pptx", md_to_pptx_bytes),
                     ("xlsx", md_to_xlsx_bytes)):
        try:
            data = fn(FULL_MD, "季度报告")
        except Exception as e:  # noqa: BLE001
            check(f"{name} 转换执行", False, str(e))
            continue
        check(f"{name} 转换输出非空", bool(data) and len(data) > 100)
        try:
            if name == "docx":
                from docx import Document
                d = Document(io.BytesIO(data))
                check("docx 可打开且有内容", len(d.paragraphs) > 0)
            elif name == "pptx":
                from pptx import Presentation
                prs = Presentation(io.BytesIO(data))
                check("pptx 可打开且分页正确（封面+2页）",
                      len(list(prs.slides)) == 3)
            else:
                from openpyxl import load_workbook
                wb = load_workbook(io.BytesIO(data))
                check("xlsx 表格 sheet 命名正确",
                      "一、业绩概览" in wb.sheetnames and "内容" in wb.sheetnames,
                      str(wb.sheetnames))
                ws = wb["一、业绩概览"]
                check("xlsx 表格数据正确",
                      ws["A1"].value == "指标" and ws["C2"].value == "112",
                      f"A1={ws['A1'].value} C2={ws['C2'].value}")
        except Exception as e:  # noqa: BLE001
            check(f"{name} 可重新打开", False, str(e))


def test_pptx_boundaries() -> None:
    from tools.doc_export import md_to_pptx_bytes
    from pptx import Presentation
    md = "\n".join(f"## 第{i}节" for i in range(45))
    prs = Presentation(io.BytesIO(md_to_pptx_bytes(md)))
    check("pptx 页数上限 40", len(list(prs.slides)) == 40)
    md2 = "## 页\n\n```py\n" + "\n".join(f"line{i}" for i in range(30)) + "\n```\n"
    prs2 = Presentation(io.BytesIO(md_to_pptx_bytes(md2)))
    check("pptx 代码块截断不报错", len(list(prs2.slides)) >= 1)
    prs3 = Presentation(io.BytesIO(md_to_pptx_bytes("")))
    check("pptx 空文档仅封面", len(list(prs3.slides)) == 1)
    md3 = "## 表\n\n| c |\n|---|\n" + "\n".join(f"| {i} |" for i in range(30)) + "\n"
    prs4 = Presentation(io.BytesIO(md_to_pptx_bytes(md3)))
    check("pptx 大表格截断不报错", len(list(prs4.slides)) >= 1)
    md4 = "# 主标题\n\n## 甲\n\n正文甲\n\n# 附录\n\n正文附录\n"
    prs5 = Presentation(io.BytesIO(md_to_pptx_bytes(md4)))
    check("pptx 多 h1 不丢内容", len(list(prs5.slides)) >= 3)


def test_xlsx_boundaries() -> None:
    from tools.doc_export import md_to_xlsx_bytes
    from openpyxl import load_workbook
    md = "| a | b |\n|---|---|\n| =SUM(A1) | +cmd |\n| -x | @y |\n"
    wb = load_workbook(io.BytesIO(md_to_xlsx_bytes(md)))
    ws = wb["表格"]
    check("xlsx 公式注入防护",
          ws["A2"].value == "'=SUM(A1)" and ws["B2"].value == "'+cmd"
          and ws["A3"].value == "'-x" and ws["B3"].value == "'@y",
          f"A2={ws['A2'].value} B2={ws['B2'].value}")
    md2 = ("## 名称:1\n\n|a|b|\n|---|---|\n|1|2|\n\n"
           "## 名称:1\n\n|a|b|\n|---|---|\n|3|4|\n")
    wb2 = load_workbook(io.BytesIO(md_to_xlsx_bytes(md2)))
    check("xlsx sheet 名清洗与冲突处理",
          "名称1" in wb2.sheetnames, str(wb2.sheetnames))
    wb3 = load_workbook(io.BytesIO(md_to_xlsx_bytes("")))
    check("xlsx 空文档有内容 sheet", wb3.sheetnames == ["内容"])
    wb4 = load_workbook(io.BytesIO(md_to_xlsx_bytes("你好\n\n世界")))
    check("xlsx 无表格走内容 sheet", "内容" in wb4.sheetnames)
    wb5 = load_workbook(io.BytesIO(md_to_xlsx_bytes("# 周报\n\n正文", "周报")))
    ws5 = wb5["内容"]
    check("xlsx title 与 h1 不重复写",
          ws5["A1"].value == "周报" and ws5["A2"].value != "周报",
          f"A1={ws5['A1'].value} A2={ws5['A2'].value}")


async def test_tool_level(tmp: Path) -> None:
    from types import SimpleNamespace
    from tools.base import ToolRegistry
    from tools.builtin import register_builtins
    reg = ToolRegistry()
    cfg = SimpleNamespace(get=lambda k, d=None: d, get_raw=lambda k, d=None: d)
    register_builtins(reg, palace=None, retriever=None, file_writer=None,
                      sandbox=None, data_dir=tmp, config=cfg)
    tool = reg.get("generate_document")
    check("generate_document 已注册", tool is not None)
    check("schema 枚举四格式",
          tool.spec.parameters["properties"]["format"]["enum"]
          == ["docx", "md", "pptx", "xlsx"])
    for fmt in ("docx", "md", "pptx", "xlsx"):
        try:
            res = await tool.run(title="冒烟", format=fmt, content=FULL_MD)
            ok = bool(res.get("download_url")) and res.get("size_bytes", 0) > 0
            check(f"工具级 {fmt} 生成", ok, str(res))
            from urllib.parse import unquote
            p = tmp / "temp" / "exports" / unquote(
                Path(res["download_url"]).name)
            check(f"工具级 {fmt} 文件落盘",
                  p.exists() and p.stat().st_size > 0)
        except Exception as e:  # noqa: BLE001
            check(f"工具级 {fmt} 生成", False, str(e))
    res = await tool.run(title="别名", format="ppt", content="x")
    check("别名 ppt→pptx", res["filename"].endswith(".pptx"))
    res = await tool.run(title="别名", format="excel", content="y")
    check("别名 excel→xlsx", res["filename"].endswith(".xlsx"))
    try:
        await tool.run(title="坏", format="pdf", content="z")
        check("非法格式被拒", False)
    except ValueError as e:
        check("非法格式被拒", "不支持的格式" in str(e))
    try:
        await tool.run(title="空", format="docx", content="  ")
        check("空内容被拒", False)
    except ValueError:
        check("空内容被拒", True)


def test_frontend_no_hardcode() -> None:
    """前端下载卡片不依赖扩展名白名单：新增格式对 UI 零改动。"""
    import re
    src = ROOT / "frontend" / "src"
    hits = []
    for p in src.rglob("*.vue"):
        text = p.read_text(encoding="utf-8")
        if re.search(r'\.docx|\.md["\']|format\s*===?\s*["\'](docx|md)["\']',
                     text):
            hits.append(str(p))
    check("前端无 docx/md 格式硬编码", not hits, "; ".join(hits))


def main() -> None:
    tmp = Path(tempfile.mkdtemp(prefix="sp_doc_export_test_"))
    test_converters()
    test_pptx_boundaries()
    test_xlsx_boundaries()
    asyncio.run(test_tool_level(tmp))
    test_frontend_no_hardcode()
    print(f"\n结果：{len(FAIL)} 项失败")
    if FAIL:
        print("失败项：", "、".join(FAIL))
        sys.exit(1)
    print("文档导出多格式测试全部通过")


if __name__ == "__main__":
    main()
